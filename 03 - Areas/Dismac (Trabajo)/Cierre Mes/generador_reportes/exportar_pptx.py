"""
Módulo para exportar reportes a PowerPoint
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os


# Colores corporativos
COLOR_FONDO = RGBColor(44, 95, 93)  # #2C5F5D
COLOR_ACENTO = RGBColor(227, 30, 36)  # #E31E24
COLOR_TEXTO_BLANCO = RGBColor(255, 255, 255)


def crear_presentacion():
    """Crea una presentación base"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    return prs


def agregar_portada(prs, mes, anio, tienda, logo_path=None):
    """Agrega diapositiva de portada"""
    slide_layout = prs.slide_layouts[6]  # Layout en blanco
    slide = prs.slides.add_slide(slide_layout)

    # Fondo
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_FONDO

    # Título principal
    left = Inches(1)
    top = Inches(2)
    width = Inches(6)
    height = Inches(1.5)

    text_box = slide.shapes.add_textbox(left, top, width, height)
    text_frame = text_box.text_frame
    text_frame.text = f"CIERRE MES\n{mes.upper()}"

    for paragraph in text_frame.paragraphs:
        paragraph.font.size = Pt(60)
        paragraph.font.bold = True
        paragraph.font.color.rgb = COLOR_TEXTO_BLANCO
        paragraph.alignment = PP_ALIGN.LEFT

    # Subtítulo
    left = Inches(1)
    top = Inches(5.5)
    width = Inches(8)
    height = Inches(1)

    text_box = slide.shapes.add_textbox(left, top, width, height)
    text_frame = text_box.text_frame
    text_frame.text = tienda

    for paragraph in text_frame.paragraphs:
        paragraph.font.size = Pt(36)
        paragraph.font.bold = False
        paragraph.font.color.rgb = COLOR_TEXTO_BLANCO
        paragraph.alignment = PP_ALIGN.CENTER

    # Logo (si existe)
    if logo_path and os.path.exists(logo_path):
        left = Inches(7.5)
        top = Inches(0.5)
        slide.shapes.add_picture(logo_path, left, top, height=Inches(1.5))

    return slide


def agregar_objetivos(prs, objetivo_usd, objetivo_bs, alcance_pct, ventas_usd, ventas_bs):
    """Agrega diapositiva de objetivos y alcance"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Fondo
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_FONDO

    # Título
    left = Inches(0.5)
    top = Inches(0.5)
    width = Inches(4)
    height = Inches(2)

    text_box = slide.shapes.add_textbox(left, top, width, height)
    text_frame = text_box.text_frame
    text_frame.text = f"OBJETIVO MES\n{objetivo_usd:,} $us\n{objetivo_bs:,} Bs"

    for paragraph in text_frame.paragraphs:
        paragraph.font.size = Pt(32)
        paragraph.font.bold = True
        paragraph.font.color.rgb = COLOR_TEXTO_BLANCO
        paragraph.alignment = PP_ALIGN.LEFT

    # Alcance
    left = Inches(5.5)
    top = Inches(1)
    width = Inches(4)
    height = Inches(2)

    text_box = slide.shapes.add_textbox(left, top, width, height)
    text_frame = text_box.text_frame
    text_frame.text = f"ALCANZE AL\n{alcance_pct}%"

    for paragraph in text_frame.paragraphs:
        paragraph.font.size = Pt(44)
        paragraph.font.bold = True
        paragraph.font.color.rgb = COLOR_TEXTO_BLANCO
        paragraph.alignment = PP_ALIGN.LEFT

    # Ventas
    left = Inches(5.5)
    top = Inches(4)
    width = Inches(4)
    height = Inches(1.5)

    text_box = slide.shapes.add_textbox(left, top, width, height)
    text_frame = text_box.text_frame
    text_frame.text = f"{ventas_usd:,} $us"

    for paragraph in text_frame.paragraphs:
        paragraph.font.size = Pt(40)
        paragraph.font.bold = True
        paragraph.font.color.rgb = COLOR_TEXTO_BLANCO
        paragraph.alignment = PP_ALIGN.LEFT

    return slide


def agregar_ticket_promedio(prs, ticket_usd, ticket_bs, contado_pct, minicuotas_pct, grafico_path=None):
    """Agrega diapositiva de ticket promedio y ventas por tipo"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Fondo
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_FONDO

    # Título izquierdo
    left = Inches(0.5)
    top = Inches(0.5)
    width = Inches(4)
    height = Inches(1)

    text_box = slide.shapes.add_textbox(left, top, width, height)
    text_frame = text_box.text_frame
    text_frame.text = "VENTAS AL:"

    for paragraph in text_frame.paragraphs:
        paragraph.font.size = Pt(36)
        paragraph.font.bold = True
        paragraph.font.color.rgb = COLOR_TEXTO_BLANCO
        paragraph.alignment = PP_ALIGN.LEFT

    # Título derecho
    left = Inches(5.5)
    top = Inches(0.5)
    width = Inches(4)
    height = Inches(1.5)

    text_box = slide.shapes.add_textbox(left, top, width, height)
    text_frame = text_box.text_frame
    text_frame.text = "TICKET\nPROMEDIO\nMES"

    for paragraph in text_frame.paragraphs:
        paragraph.font.size = Pt(28)
        paragraph.font.bold = True
        paragraph.font.color.rgb = COLOR_TEXTO_BLANCO
        paragraph.alignment = PP_ALIGN.LEFT

    # Ticket promedio
    left = Inches(5.5)
    top = Inches(2.5)
    width = Inches(4)
    height = Inches(1.5)

    text_box = slide.shapes.add_textbox(left, top, width, height)
    text_frame = text_box.text_frame
    text_frame.text = f"{ticket_usd} $us"

    for paragraph in text_frame.paragraphs:
        paragraph.font.size = Pt(40)
        paragraph.font.bold = True
        paragraph.font.color.rgb = COLOR_TEXTO_BLANCO
        paragraph.alignment = PP_ALIGN.LEFT

    # Gráfico de dona (si existe)
    if grafico_path and os.path.exists(grafico_path):
        left = Inches(0.5)
        top = Inches(2)
        slide.shapes.add_picture(grafico_path, left, top, width=Inches(4.5))

    return slide


def agregar_grafico_categorias(prs, titulo, grafico_path):
    """Agrega diapositiva con gráfico de categorías"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Fondo
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_FONDO

    # Título
    left = Inches(0.5)
    top = Inches(0.3)
    width = Inches(9)
    height = Inches(0.8)

    text_box = slide.shapes.add_textbox(left, top, width, height)
    text_frame = text_box.text_frame
    text_frame.text = titulo

    for paragraph in text_frame.paragraphs:
        paragraph.font.size = Pt(32)
        paragraph.font.bold = True
        paragraph.font.color.rgb = COLOR_TEXTO_BLANCO
        paragraph.alignment = PP_ALIGN.CENTER

    # Gráfico
    if grafico_path and os.path.exists(grafico_path):
        left = Inches(0.3)
        top = Inches(1.3)
        slide.shapes.add_picture(grafico_path, left, top, width=Inches(9.4))

    return slide


def agregar_mystery_shopper(prs, indice_general, atencion_vendedor, puntos_mejorar):
    """
    Agrega diapositiva de Mystery Shopper
    puntos_mejorar: lista de diccionarios con 'descripcion' y 'porcentaje'
    """
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Fondo
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_FONDO

    # Título
    left = Inches(0.5)
    top = Inches(0.3)
    width = Inches(4)
    height = Inches(1.5)

    text_box = slide.shapes.add_textbox(left, top, width, height)
    text_frame = text_box.text_frame
    text_frame.text = "MISTERY\nSHOPPER"

    for paragraph in text_frame.paragraphs:
        paragraph.font.size = Pt(40)
        paragraph.font.bold = True
        paragraph.font.color.rgb = COLOR_TEXTO_BLANCO
        paragraph.alignment = PP_ALIGN.LEFT

    # Subtítulo
    left = Inches(0.5)
    top = Inches(2.5)
    width = Inches(4)
    height = Inches(1)

    text_box = slide.shapes.add_textbox(left, top, width, height)
    text_frame = text_box.text_frame
    text_frame.text = "PUNTOS\nA\nMEJORAR"

    for paragraph in text_frame.paragraphs:
        paragraph.font.size = Pt(32)
        paragraph.font.bold = True
        paragraph.font.color.rgb = COLOR_TEXTO_BLANCO
        paragraph.alignment = PP_ALIGN.LEFT

    # Indicadores
    left = Inches(5)
    top = Inches(0.5)
    width = Inches(4.5)
    height = Inches(1)

    text_box = slide.shapes.add_textbox(left, top, width, height)
    text_frame = text_box.text_frame
    text_frame.text = f"ÍNDICE GENERAL: {indice_general}%\nATENCIÓN DEL VENDEDOR: {atencion_vendedor}%"

    for paragraph in text_frame.paragraphs:
        paragraph.font.size = Pt(14)
        paragraph.font.bold = True
        paragraph.font.color.rgb = COLOR_TEXTO_BLANCO
        paragraph.alignment = PP_ALIGN.LEFT

    # Puntos a mejorar (tabla simple)
    left = Inches(5)
    top = Inches(1.8)
    width = Inches(4.5)

    y_offset = 0
    for punto in puntos_mejorar[:5]:  # Máximo 5 puntos
        text_box = slide.shapes.add_textbox(left, top + Inches(y_offset), width, Inches(0.5))
        text_frame = text_box.text_frame
        text_frame.text = f"• {punto['descripcion']}: {punto['porcentaje']}%"
        text_frame.word_wrap = True

        for paragraph in text_frame.paragraphs:
            paragraph.font.size = Pt(11)
            paragraph.font.color.rgb = COLOR_TEXTO_BLANCO
            paragraph.alignment = PP_ALIGN.LEFT

        y_offset += 0.6

    return slide


def agregar_gracias(prs, logo_path=None):
    """Agrega diapositiva final de agradecimiento"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Fondo
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_FONDO

    # Texto
    left = Inches(2)
    top = Inches(2.5)
    width = Inches(6)
    height = Inches(2)

    text_box = slide.shapes.add_textbox(left, top, width, height)
    text_frame = text_box.text_frame
    text_frame.text = "GRACIAS"

    for paragraph in text_frame.paragraphs:
        paragraph.font.size = Pt(72)
        paragraph.font.bold = True
        paragraph.font.color.rgb = COLOR_TEXTO_BLANCO
        paragraph.alignment = PP_ALIGN.CENTER

    # Logo (si existe)
    if logo_path and os.path.exists(logo_path):
        left = Inches(3.5)
        top = Inches(4.5)
        slide.shapes.add_picture(logo_path, left, top, height=Inches(2))

    return slide


def guardar_presentacion(prs, ruta_salida):
    """Guarda la presentación en la ruta especificada"""
    prs.save(ruta_salida)
    return ruta_salida
